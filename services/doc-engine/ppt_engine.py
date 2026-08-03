"""Markdown outline → editable PPTX for Free myself PPT Studio."""

from __future__ import annotations

import base64
import io
import re
from dataclasses import dataclass, field
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# 16:9 widescreen
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)


@dataclass
class Theme:
    id: str
    label: str
    hint: str
    bg: str
    accent: str
    title_color: str
    body_color: str
    muted_color: str
    title_font: str
    body_font: str
    accent_bar: bool = True


THEMES: dict[str, Theme] = {
    "business-light": Theme(
        id="business-light",
        label="浅色商务",
        hint="蓝灰配色，适合周会与汇报",
        bg="#F7F9FC",
        accent="#2563EB",
        title_color="#0F172A",
        body_color="#334155",
        muted_color="#64748B",
        title_font="微软雅黑",
        body_font="微软雅黑",
    ),
    "academic-clean": Theme(
        id="academic-clean",
        label="简洁学术",
        hint="白底深字，适合课程与答辩",
        bg="#FFFFFF",
        accent="#1E3A5F",
        title_color="#111827",
        body_color="#374151",
        muted_color="#6B7280",
        title_font="宋体",
        body_font="宋体",
    ),
    "minimal-ink": Theme(
        id="minimal-ink",
        label="极简墨色",
        hint="柔灰背景，强调排版层次",
        bg="#F4F4F5",
        accent="#18181B",
        title_color="#18181B",
        body_color="#3F3F46",
        muted_color="#71717A",
        title_font="微软雅黑",
        body_font="微软雅黑",
        accent_bar=False,
    ),
}


@dataclass
class SlidePlan:
    kind: str  # cover | content | section | closing
    title: str
    bullets: list[str] = field(default_factory=list)
    subtitle: str = ""


def theme_catalog() -> dict[str, Any]:
    return {
        "themes": [
            {
                "id": t.id,
                "label": t.label,
                "hint": t.hint,
            }
            for t in THEMES.values()
        ],
        "aspect": "16:9",
    }


def resolve_theme(spec: dict[str, Any] | None) -> Theme:
    theme_id = "business-light"
    if isinstance(spec, dict):
        raw = spec.get("themeId") or spec.get("theme") or theme_id
        if isinstance(raw, str) and raw in THEMES:
            theme_id = raw
    return THEMES[theme_id]


def _hex_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        h = "000000"
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _strip_inline(md: str) -> str:
    text = md.strip()
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", text)
    return text.strip()


def parse_markdown_outline(markdown: str) -> list[SlidePlan]:
    """Parse `#` cover, `##` slides, `-` bullets, optional `---` breaks."""
    lines = markdown.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    slides: list[SlidePlan] = []
    cover_title = ""
    cover_subtitle = ""
    current: SlidePlan | None = None
    saw_h1 = False

    def flush() -> None:
        nonlocal current
        if current and (current.title or current.bullets):
            slides.append(current)
        current = None

    for raw in lines:
        line = raw.rstrip()
        stripped = line.strip()
        if not stripped:
            continue
        if stripped == "---":
            flush()
            continue

        h1 = re.match(r"^#\s+(.+)$", stripped)
        if h1:
            flush()
            if not saw_h1:
                cover_title = _strip_inline(h1.group(1))
                saw_h1 = True
            else:
                current = SlidePlan(kind="section", title=_strip_inline(h1.group(1)))
            continue

        h2 = re.match(r"^##\s+(.+)$", stripped)
        if h2:
            flush()
            title = _strip_inline(h2.group(1))
            kind = "closing" if _is_closing_title(title) else "content"
            current = SlidePlan(kind=kind, title=title)
            continue

        h3 = re.match(r"^###\s+(.+)$", stripped)
        if h3:
            text = _strip_inline(h3.group(1))
            if current is None and saw_h1 and not cover_subtitle:
                cover_subtitle = text
            elif current is not None:
                current.bullets.append(text)
            else:
                current = SlidePlan(kind="content", title=text)
            continue

        bullet = re.match(r"^[-*+]\s+(.+)$", stripped) or re.match(
            r"^\d+\.\s+(.+)$", stripped
        )
        if bullet:
            text = _strip_inline(bullet.group(1))
            if current is None:
                if saw_h1 and not cover_subtitle and not slides:
                    cover_subtitle = text
                else:
                    current = SlidePlan(kind="content", title="要点", bullets=[text])
            else:
                current.bullets.append(text)
            continue

        # Plain paragraph
        text = _strip_inline(stripped)
        if not text:
            continue
        if current is None:
            if saw_h1 and not cover_subtitle:
                cover_subtitle = text
            elif not saw_h1 and not cover_title:
                cover_title = text
                saw_h1 = True
            else:
                current = SlidePlan(kind="content", title=text)
        else:
            current.bullets.append(text)

    flush()

    if not cover_title and slides:
        first = slides[0]
        if first.kind == "content" and not first.bullets:
            cover_title = first.title
            slides = slides[1:]
        else:
            cover_title = first.title

    if not cover_title:
        cover_title = "演示文稿"

    result = [
        SlidePlan(kind="cover", title=cover_title, subtitle=cover_subtitle),
        *slides,
    ]

    if not any(s.kind == "closing" for s in result) and len(result) >= 2:
        # Keep decks short for daily use — no forced closing slide.
        pass

    return result


def _is_closing_title(title: str) -> bool:
    t = title.strip().lower()
    keys = ("谢谢", "感谢", "q&a", "qa", "问答", "结束", "thank")
    return any(k in t for k in keys)


def _set_run_font(run: Any, font_name: str, size_pt: float, color: str, bold: bool = False) -> None:
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _hex_rgb(color)
    # East Asian font hint for PowerPoint/WPS
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            from lxml import etree

            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", font_name)


def _fill_bg(slide: Any, theme: Theme) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_rgb(theme.bg)


def _add_accent_bar(slide: Any, theme: Theme) -> None:
    if not theme.accent_bar:
        return
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(0.18),
        _SLIDE_H,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_rgb(theme.accent)
    shape.line.fill.background()


def _add_textbox(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_name: str,
    size_pt: float,
    color: str,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, font_name, size_pt, color, bold=bold)


def _render_cover(slide: Any, plan: SlidePlan, theme: Theme) -> None:
    _fill_bg(slide, theme)
    # Bottom accent band
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(6.6),
        _SLIDE_W,
        Inches(0.9),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = _hex_rgb(theme.accent)
    band.line.fill.background()

    _add_textbox(
        slide,
        1.0,
        2.2,
        11.0,
        1.6,
        plan.title,
        font_name=theme.title_font,
        size_pt=40,
        color=theme.title_color,
        bold=True,
        align=PP_ALIGN.LEFT,
    )
    if plan.subtitle:
        _add_textbox(
            slide,
            1.0,
            4.0,
            11.0,
            0.8,
            plan.subtitle,
            font_name=theme.body_font,
            size_pt=18,
            color=theme.muted_color,
            bold=False,
        )


def _render_content(slide: Any, plan: SlidePlan, theme: Theme) -> None:
    _fill_bg(slide, theme)
    _add_accent_bar(slide, theme)
    _add_textbox(
        slide,
        0.7,
        0.45,
        11.8,
        1.0,
        plan.title,
        font_name=theme.title_font,
        size_pt=28,
        color=theme.title_color,
        bold=True,
    )
    # Accent underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7),
        Inches(1.45),
        Inches(1.6),
        Inches(0.06),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _hex_rgb(theme.accent)
    line.line.fill.background()

    if not plan.bullets:
        _add_textbox(
            slide,
            0.7,
            2.0,
            11.5,
            4.5,
            "（可在此页补充要点）",
            font_name=theme.body_font,
            size_pt=16,
            color=theme.muted_color,
        )
        return

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(11.5), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(plan.bullets[:12]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = f"•  {bullet}"
        _set_run_font(run, theme.body_font, 18, theme.body_color, bold=False)


def _render_closing(slide: Any, plan: SlidePlan, theme: Theme) -> None:
    _fill_bg(slide, theme)
    _add_textbox(
        slide,
        1.0,
        2.8,
        11.3,
        1.4,
        plan.title,
        font_name=theme.title_font,
        size_pt=36,
        color=theme.title_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    if plan.bullets:
        _add_textbox(
            slide,
            1.0,
            4.4,
            11.3,
            1.0,
            "  ·  ".join(plan.bullets[:4]),
            font_name=theme.body_font,
            size_pt=16,
            color=theme.muted_color,
            align=PP_ALIGN.CENTER,
        )


def markdown_to_pptx(markdown: str, spec: dict[str, Any] | None = None) -> tuple[bytes, dict[str, Any]]:
    theme = resolve_theme(spec)
    plans = parse_markdown_outline(markdown)
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    blank = prs.slide_layouts[6]  # blank

    kind_counts: dict[str, int] = {}
    for plan in plans:
        slide = prs.slides.add_slide(blank)
        if plan.kind == "cover":
            _render_cover(slide, plan, theme)
        elif plan.kind == "closing":
            _render_closing(slide, plan, theme)
        else:
            _render_content(slide, plan, theme)
        kind_counts[plan.kind] = kind_counts.get(plan.kind, 0) + 1

    buf = io.BytesIO()
    prs.save(buf)
    out = buf.getvalue()
    summary = {
        "slideCount": len(plans),
        "themeId": theme.id,
        "kinds": kind_counts,
        "titleFont": theme.title_font,
        "bodyFont": theme.body_font,
    }
    return out, summary


def bytes_to_b64(data: bytes) -> str:
    return base64.b64encode(data).decode("ascii")
